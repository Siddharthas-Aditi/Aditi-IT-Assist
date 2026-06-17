"""Text extraction from supported document formats.

Supported formats:
- ``.docx`` — python-docx
- ``.pdf``  — pymupdf (fitz)
- ``.pptx`` — python-pptx
- ``.txt``  — plain UTF-8 read
- ``.md``   — plain UTF-8 read (markdown preserved as-is)

All extractors return an ``ExtractionResult`` containing the raw text and
optional structural metadata (e.g. slide titles, heading levels).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


# ── Result type ───────────────────────────────────────────────────────────────

@dataclass
class ExtractionResult:
    """Outcome of text extraction from a single document."""

    raw_text: str
    page_count: int = 0
    word_count: int = 0
    structural_metadata: dict = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)

    def __post_init__(self) -> None:
        if not self.word_count:
            self.word_count = len(self.raw_text.split())


# ── Dispatcher ────────────────────────────────────────────────────────────────

def extract_text(file_path: str | Path) -> ExtractionResult:
    """Extract raw text from *file_path*.

    The format is detected from the file extension.  Raises ``ValueError``
    for unsupported extensions and ``RuntimeError`` if extraction fails
    completely.
    """
    path = Path(file_path)
    ext = path.suffix.lower().lstrip(".")

    extractors = {
        "docx": _extract_docx,
        "pdf": _extract_pdf,
        "pptx": _extract_pptx,
        "txt": _extract_plaintext,
        "md": _extract_plaintext,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        msg = f"Unsupported file extension '.{ext}'.  Allowed: {list(extractors)}"
        raise ValueError(msg)

    try:
        return extractor(path)
    except Exception as exc:
        logger.exception("Extraction failed for %s", path)
        raise RuntimeError(f"Text extraction failed: {exc}") from exc


# ── Format-specific extractors ────────────────────────────────────────────────

def _extract_docx(path: Path) -> ExtractionResult:
    """Extract text from a .docx file using python-docx."""
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError as e:
        msg = "python-docx is required for .docx extraction.  Run: pip install python-docx"
        raise ImportError(msg) from e

    doc = Document(str(path))
    paragraphs: list[str] = []
    heading_styles: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Capture heading level for structural metadata
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            heading_styles.append(para.style.name)
        paragraphs.append(text)

    raw_text = "\n".join(paragraphs)
    return ExtractionResult(
        raw_text=raw_text,
        page_count=0,  # DOCX has no native page concept without rendering
        structural_metadata={"heading_styles": heading_styles},
    )


def _extract_pdf(path: Path) -> ExtractionResult:
    """Extract text from a .pdf file using pymupdf (fitz)."""
    try:
        import fitz  # type: ignore[import-untyped]  # pymupdf
    except ImportError as e:
        msg = "pymupdf is required for .pdf extraction.  Run: pip install pymupdf"
        raise ImportError(msg) from e

    doc = fitz.open(str(path))
    pages: list[str] = []
    warnings: list[str] = []

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("blocks")  # list of (x0, y0, x1, y1, text, ...)
        page_lines: list[str] = []
        for block in blocks:
            block_text = block[4].strip() if len(block) > 4 else ""
            if block_text:
                page_lines.append(block_text)
        if not page_lines:
            warnings.append(f"Page {page_num} yielded no text (possibly scanned image).")
        pages.append("\n".join(page_lines))

    doc.close()
    raw_text = "\n\n".join(p for p in pages if p)
    return ExtractionResult(
        raw_text=raw_text,
        page_count=len(pages),
        warnings=warnings,
    )


def _extract_pptx(path: Path) -> ExtractionResult:
    """Extract text from a .pptx file using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as e:
        msg = "python-pptx is required for .pptx extraction.  Run: pip install python-pptx"
        raise ImportError(msg) from e

    prs = Presentation(str(path))
    slides_text: list[str] = []
    slide_titles: list[str] = []

    for slide in prs.slides:
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            title_candidate = slide_lines[0]
            slide_titles.append(title_candidate)
            slides_text.append("\n".join(slide_lines))

    raw_text = "\n\n".join(slides_text)
    return ExtractionResult(
        raw_text=raw_text,
        page_count=len(prs.slides),
        structural_metadata={"slide_titles": slide_titles},
    )


def _extract_plaintext(path: Path) -> ExtractionResult:
    """Read a plain .txt or .md file."""
    raw_text = path.read_text(encoding="utf-8", errors="replace")
    return ExtractionResult(raw_text=raw_text)
